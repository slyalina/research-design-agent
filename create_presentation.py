#!/usr/bin/env python3
"""
Generate PowerPoint presentation for video submission.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Paths to generated images
HOME_DIR = os.path.expanduser("~")
ARTIFACT_DIR = f"{HOME_DIR}/.gemini/antigravity/brain/913cd6e3-381a-4636-98cc-1625c8eb18e1"
IMAGES = {
    'title': f"{ARTIFACT_DIR}/title_slide_background_1764615597861.png",
    'problem': f"{ARTIFACT_DIR}/problem_statement_visual_1764615618469.png",
    'architecture': f"{ARTIFACT_DIR}/multi_agent_architecture_1764615633572.png",
    'demo': f"{ARTIFACT_DIR}/demo_power_analysis_1764615653899.png",
    'tech': f"{ARTIFACT_DIR}/tech_stack_visual_1764615668614.png",
}

def create_presentation():
    """Create the PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background image
    if os.path.exists(IMAGES['title']):
        slide.shapes.add_picture(IMAGES['title'], 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    # Add title text box
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Bioinformatics Research Design Agent"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    left = Inches(1)
    top = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1))
    tf = txBox.text_frame
    tf.text = "Your AI-Powered Research Planning Partner"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add tech stack
    top = Inches(5.5)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = txBox.text_frame
    tf.text = "Built with Google ADK | Gemini 2.5 | Multi-Agent Architecture"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(200, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "The Problem 🎯"
    
    # Add problem image
    if os.path.exists(IMAGES['problem']):
        left = Inches(0.5)
        top = Inches(1.5)
        slide.shapes.add_picture(IMAGES['problem'], left, top, width=Inches(9))
    
    # Add text content
    left = Inches(0.5)
    top = Inches(4.5)
    width = Inches(9)
    height = Inches(2.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "❌ Underpowered Studies: Improper sample size → wasted resources & failed trials"
    p.font.size = Pt(16)
    p.space_before = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "❌ Literature Overload: Days searching for effect sizes & parameters"
    p.font.size = Pt(16)
    p.space_before = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "❌ Statistical Rigor Gap: Complex analyses require specialized expertise"
    p.font.size = Pt(16)
    p.space_before = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "💰 Failed trials waste $800M-$1.4B per drug | 65% of studies not reproducible"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 0, 0)
    p.space_before = Pt(12)
    
    # Slide 3: Why Agents?
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Why Agents? 🤖"
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Multi-Agent Orchestration = Specialized Expertise"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Study design requires multiple types of expertise:"
    p.font.size = Pt(18)
    p.space_after = Pt(6)
    
    for item in ["✅ Literature expertise (search & extraction)",
                 "✅ Statistical expertise (power analysis, simulations)",
                 "✅ Domain expertise (biomarkers, protocols)",
                 "✅ Critical review (bias detection, rigor)"]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.level = 1
        p.space_before = Pt(3)
    
    p = tf.add_paragraph()
    p.text = "Why Not a Single LLM?"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "❌ Hallucinated statistics | ❌ No tool execution | ❌ Can't handle multi-step workflows"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(150, 0, 0)
    
    p = tf.add_paragraph()
    p.text = "✅ Agents provide grounded, reproducible, code-backed analysis"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 100, 0)
    p.space_before = Pt(6)
    
    # Slide 4: Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Architecture Overview 🏗️"
    
    if os.path.exists(IMAGES['architecture']):
        left = Inches(0.5)
        top = Inches(1.5)
        slide.shapes.add_picture(IMAGES['architecture'], left, top, width=Inches(9))
    
    left = Inches(0.5)
    top = Inches(5)
    width = Inches(9)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "6 Specialized Agents: Literature • Power Analysis • Biomarker Data • Microbiome • Proposal • Criticism"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Key ADK Concepts: Sub-agent delegation | Tool integration (R, MCP) | Human-in-the-loop"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_before = Pt(8)
    
    # Slide 5: Demo - Power Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Demo: Power Analysis Flow 💪"
    
    if os.path.exists(IMAGES['demo']):
        left = Inches(0.5)
        top = Inches(1.5)
        slide.shapes.add_picture(IMAGES['demo'], left, top, width=Inches(9))
    
    left = Inches(0.5)
    top = Inches(5.5)
    width = Inches(9)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = "User: 'Need sample size for diabetes trial, 0.8% HbA1c reduction, 80% power'"
    p.font.size = Pt(14)
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "Result: ✅ n=64 per group (128 total) | ✅ Generated R code | ✅ Literature citations"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 100, 0)
    p.space_before = Pt(8)
    
    # Slide 6: End-to-End Workflow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Demo: End-to-End Workflow 🔄"
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Input: 'Study gut microbiome biomarkers for Crohn's disease prognosis'"
    p.font.size = Pt(16)
    p.font.italic = True
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Agent Workflow:"
    p.font.size = Pt(18)
    p.font.bold = True
    
    workflow_steps = [
        "1. Literature Agent → Finds microbiome-Crohn's papers",
        "2. Biomarker Data Agent → Searches SRA for public datasets",
        "3. Power Analysis Agent → Calculates sample size from literature",
        "4. Microbiome Tool Runner → Processes sample data",
        "5. Proposal Agent → Drafts formal research proposal",
        "6. Criticism Agent → Reviews for bias, multiple testing, confounding"
    ]
    
    for step in workflow_steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(14)
        p.level = 1
        p.space_before = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Output: 📄 Complete proposal + 📊 Processed data + 💾 Reproducible R scripts"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 100, 0)
    p.space_before = Pt(12)
    
    # Slide 7: The Build
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Build ⚙️"
    
    if os.path.exists(IMAGES['tech']):
        left = Inches(0.5)
        top = Inches(1.5)
        slide.shapes.add_picture(IMAGES['tech'], left, top, width=Inches(9))
    
    left = Inches(0.5)
    top = Inches(5)
    width = Inches(9)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "AI: Google ADK + Gemini 2.5 Flash + MCP"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Stats: R (pwr, lme4, simr, survival) + Custom execution tool"
    p.font.size = Pt(14)
    p.space_before = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Data: SRA/ENA, GEO, Cell x Gene + Bioinformatics pipelines"
    p.font.size = Pt(14)
    p.space_before = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Deployment: Docker + Cloud Run ready"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_before = Pt(8)
    
    # Slide 8: Impact & Value
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Impact & Value 🎯"
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "What This Agent Delivers"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_after = Pt(12)
    
    benefits = [
        "⏰ Time Savings: Days → Minutes for study design",
        "📊 Statistical Rigor: Code-backed calculations, not guesses",
        "📚 Literature Grounding: Automated parameter extraction",
        "🔄 Reproducibility: All analyses saved as executable code"
    ]
    
    for benefit in benefits:
        p = tf.add_paragraph()
        p.text = benefit
        p.font.size = Pt(16)
        p.space_before = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Innovation Highlights"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(18)
    
    innovations = [
        "Hybrid AI + R approach for grounded statistics",
        "Criticism agent for methodological quality control",
        "Full pipeline: literature → analysis → critique"
    ]
    
    for innovation in innovations:
        p = tf.add_paragraph()
        p.text = "• " + innovation
        p.font.size = Pt(14)
        p.level = 1
        p.space_before = Pt(4)
    
    # Save presentation
    output_path = "Video_Submission_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
